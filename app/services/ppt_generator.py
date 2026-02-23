from pptx import Presentation
from pptx.util import Inches
from typing import List

def create_presentation(slides_data: List[dict], filename: str):
    """Створює презентацію .pptx зі списку слайдів."""
    prs = Presentation()
    for slide_info in slides_data:
        slide_layout = prs.slide_layouts[1]  # заголовок + текст
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_info.get('title', '')
        # Можна додати форматування: замінити маркери списку
        text = slide_info.get('content', '')
        content.text = text

    prs.save(filename)
    return filename
